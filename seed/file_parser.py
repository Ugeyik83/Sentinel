"""
seed/file_parser.py
Evrensel dosya parse — PDF/XLSX/DOCX/PPTX/CSV/TXT/JSON
"""

import json
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_file(file_path: str) -> str:
    path = Path(file_path)
    ext = path.suffix.lower()
    parsers = {
        ".pdf": _parse_pdf,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".docx": _parse_docx,
        ".csv": _parse_csv,
        ".txt": _parse_text,
        ".md": _parse_text,
        ".json": _parse_json,
        ".pptx": _parse_pptx,
    }
    parser = parsers.get(ext, _parse_text)
    try:
        text = parser(str(path))
        logger.info(f"Parse edildi: {path.name} — {len(text):,} karakter")
        return text
    except Exception as e:
        logger.error(f"Parse hatası [{path.name}]: {e}")
        return f"[Parse hatası: {e}]"


def _parse_pdf(path):
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n\n".join(p.extract_text() or "" for p in reader.pages)
    except ImportError:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n\n".join(p.extract_text() or "" for p in pdf.pages)


def _parse_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _parse_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_csv(path):
    import pandas as pd
    df = pd.read_csv(path, nrows=1000)
    return df.to_string(index=False)


def _parse_text(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _parse_json(path):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)


def _parse_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slayt {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n\n".join(parts)
